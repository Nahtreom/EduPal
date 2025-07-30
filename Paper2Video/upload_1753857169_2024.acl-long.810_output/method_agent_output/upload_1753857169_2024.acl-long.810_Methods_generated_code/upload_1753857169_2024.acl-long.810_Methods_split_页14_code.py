import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "LLM Hallucinations in Software Development"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "LLM hallucinations manifest when models generate outputs that are nonsensical,",
    "factually incorrect, or inaccurate (Dhuliawala et al., 2023; Zhang et al., 2023b).",
    "This issue is particularly concerning in software development, where programming",
    "languages demand precise syntax—the absence of even a single line can lead to",
    "system failure. We have observed that LLMs often produce coding hallucinations,",
    "which encompass potential issues like incomplete implementations, unexecutable",
    "code, and inconsistencies that don't meet requirements."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(10)

# Add image if path is provided
# image_path = "path/to/image.png"  # Replace with actual image path if available
# if os.path.exists(image_path):
#     slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))