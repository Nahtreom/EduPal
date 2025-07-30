import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Enhancing Quality with ChatDev"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "To enhance the quality and reduce human intervention, ChatDev implements prompt engineering that only takes place at the start of each subtask round.",
    "As soon as the communication phase begins, the instructor and the assistant will communicate with each other in an automated loop, continuing this exchange until the task concludes.",
    "However, simply exchanging responses cannot achieve effective multi-round task-oriented communication, since it inevitably faces significant challenges including role flipping, instruction repeating, and fake replies."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)

# Add a placeholder for an image (if needed)
# image_path = "path/to/image.png"  # Uncomment and set the correct path if an image is available
# slide.shapes.add_picture(image_path, Inches(6), Inches(3), width=Inches(2), height=Inches(2))