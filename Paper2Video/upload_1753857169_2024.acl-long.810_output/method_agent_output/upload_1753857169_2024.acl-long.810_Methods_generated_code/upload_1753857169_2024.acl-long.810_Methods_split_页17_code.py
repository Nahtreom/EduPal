import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Optimization Process"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "After the instructor provides a specific modification suggestion, the assistant proceeds to perform precise optimization:"
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Mathematical expression
math_expr = content_frame.add_paragraph()
math_expr.text = "<math expression> <langle \\mathcal { T } \\to \\mathcal { A } , \\langle \\mathcal { A } \\to \\mathcal { I } , \\mathcal { I } \\sim \\mathcal { A } \\rangle _ { \\odot } , \\mathcal { A } \\sim \\mathcal { I } \\rangle _ { \\odot } </math expression>"
math_expr.font.size = Pt(18)
math_expr.font.color.rgb = RGBColor(0, 0, 0)

# Additional content
additional_content = content_frame.add_paragraph()
additional_content.text = "Since this mechanism tackles one concrete issue at a time, it requires multiple rounds of communication to optimize various potential problems. The communication pattern instructs agents on how to communicate, enabling finer-grained information exchange for effective solution optimization, which practically aids in reducing coding hallucinations."
additional_content.font.size = Pt(18)
additional_content.font.color.rgb = RGBColor(0, 0, 0)

# Add image if path is provided
image_path = "path/to/image.png"  # Replace with actual image path if available
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))