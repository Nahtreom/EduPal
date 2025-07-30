import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Bridging Solutions in Software Development"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = ("Subsequently, the solutions from previous tasks serve as bridges to the next phase, "
                 "allowing a smooth transition between subtasks. This approach continues until all subtasks "
                 "are completed. It's worth noting that the conceptually simple but empirically powerful "
                 "chain-style structure guides agents on what to communicate, fostering cooperation and "
                 "smoothly linking natural- and programming-language subtasks. It also offers a transparent "
                 "view of the entire software development process, allowing for the examination of intermediate "
                 "solutions and assisting in identifying possible problems.")
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)
content.alignment = PP_ALIGN.LEFT

# Add a table if needed (not applicable here, but included for completeness)
# table_data = [['Header1', 'Header2'], ['Row1Col1', 'Row1Col2']]
# rows, cols = len(table_data), len(table_data[0])
# left = Inches(1)
# top = Inches(3.5)
# width = Inches(8)
# height = Inches(1)
# table = slide.shapes.add_table(rows, cols, left, top, width, height).table
# for r in range(rows):
#     for c in range(cols):
#         table.cell(r, c).text = table_data[r][c]

# Add an image if needed (not applicable here, but included for completeness)
# img_path = 'path/to/image.png'
# slide.shapes.add_picture(img_path, Inches(6), Inches(4), width=Inches(2), height=Inches(2))