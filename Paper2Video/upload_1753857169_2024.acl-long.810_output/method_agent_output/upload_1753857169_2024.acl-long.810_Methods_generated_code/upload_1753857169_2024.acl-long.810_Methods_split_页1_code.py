import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Introduction to ChatDev"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "We introduce ChatDev, a chat-powered software development framework that integrates multiple 'software agents' with various social roles (e.g., requirements analysts, professional programmers and test engineers) collaborating in the core phases of the software life cycle.",
    "Technically, to facilitate cooperative communication, ChatDev introduces chat chain to further break down each phase into smaller and manageable subtasks, which guides multi-turn communications between different roles to propose and validate solutions for each subtask.",
    "In addition, to alleviate unexpected hallucinations, a communicative pattern named communicative dehallucination is devised, wherein agents request more detailed information before responding directly and then continue the next round of communication based on these details."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)

# Image
img_path = "path/to/your/image.png"  # Replace with actual image path
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(5.5), Inches(3), width=Inches(2), height=Inches(2))