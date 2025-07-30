import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Instructor and Assistant Instantiation"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "Then, an instructor T and an assistant A are instantiated by hypnotizing LLMs via P_I and P_A."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "T = ρ(LLM, P_I), A = ρ(LLM, P_A)"
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "where ρ is the role customization operation, implemented via system message assignment."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Add image if path is provided
image_path = "path/to/image.png"  # Replace with actual image path if available
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))