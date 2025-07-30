import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Memory Segmentation in LLMs"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "Note that the limited context length of common LLMs typically restricts the ability to maintain a complete communication history among all agents and phases.",
    "To tackle this issue, based on the nature of the chat chain, we accordingly segment the agents’ context memories based on their sequential phases, resulting in two functionally distinct types of memory:",
    "1. Short-term memory: utilized to sustain the continuity of the dialogue within a single phase.",
    "2. Long-term memory: leveraged to preserve contextual awareness across different phases."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)

# Optional: Add an image if a path is provided
# image_path = "path/to/image.png"  # Replace with actual image path if available
# if os.path.exists(image_path):
#     slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))