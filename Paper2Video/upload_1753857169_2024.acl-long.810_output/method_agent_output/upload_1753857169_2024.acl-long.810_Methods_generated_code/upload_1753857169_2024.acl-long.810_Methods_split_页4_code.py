import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Agentic Workflow Task Solving Process"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "The entire task solving process along the agentic workflow can be formulated as:"
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Mathematical formula
formula_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
formula_frame = formula_box.text_frame
formula_frame.word_wrap = True

formula = formula_frame.add_paragraph()
formula.text = "C = < P^1, P^2, ..., P^{|C|} >\nP^i = < T^1, T^2, ..., T^{|P^i|} >\nT^j = τ(C(Z, A))\nC(Z, A) = < Z A, A Z >_⊙"
formula.font.size = Pt(16)
formula.font.color.rgb = RGBColor(0, 0, 0)

# Additional content
additional_content = content_frame.add_paragraph()
additional_content.text = "The dual-agent communication design simplifies communications by avoiding complex multi-agent topologies, effectively streamlining the consensus reaching process (Yin et al., 2023; Chen et al., 2023b)."
additional_content.font.size = Pt(18)
additional_content.font.color.rgb = RGBColor(0, 0, 0)

# Optional image (if path is provided)
# image_path = "path/to/image.png"  # Replace with actual image path if available
# if os.path.exists(image_path):
#     slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(3), height=Inches(2))