import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Communication Patterns"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# First paragraph
p1 = content_frame.add_paragraph()
p1.text = "Specifically, a vanilla communication pattern between the assistant and the instructor follows a straightforward instruction-response format:"
p1.font.size = Pt(20)
p1.font.color.rgb = RGBColor(0, 0, 0)

# Math formula
p2 = content_frame.add_paragraph()
p2.text = "<math> <langle \\mathcal { T } \\to \\mathcal { A } , \\mathcal { A } \\{ \\mathcal { A } \\} \\emptyset </math>"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(0, 0, 0)

# Second paragraph
p3 = content_frame.add_paragraph()
p3.text = "In contrast, our communicative dehallucination mechanism features a deliberate 'role reversal', where the assistant takes on an instructor-like role, proactively seeking more specific information (e.g., the precise name of an external dependency and its related class) before delivering a conclusive response."
p3.font.size = Pt(20)
p3.font.color.rgb = RGBColor(0, 0, 0)

# Optional image (if path is provided)
image_path = "path/to/image.png"  # Replace with actual image path if available
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5.5), Inches(3), width=Inches(2), height=Inches(2))