import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Agent Collaboration in Subtasks"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "In every subtask, two agents, each with their own specialized roles (e.g., a reviewer skilled at identifying endless loops and a programmer adept in GUI design), perform the functions of an instructor (T) and an assistant (A).",
    "The instructor agent initiates instructions, instructing ( ) the discourse toward the completion of the subtask, while the assistant agent adheres to these instructions and responds with ( ) appropriate solutions.",
    "They engage in a multi-turn dialogue (C), working cooperatively until they achieve consensus, extracting (τ) solutions that can range from the text (e.g. defining a software function point) to code (e.g., creating the initial version of source code), ultimately leading to the completion of the subtask."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)

# Add image if path is provided
image_path = "path/to/image.png"  # Replace with actual image path if available
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5.5), Inches(3), width=Inches(2), height=Inches(2))