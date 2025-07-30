import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Communicative Dehallucination"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "Coding hallucinations frequently appear when the assistant struggles to precisely follow instructions, often due to the vagueness and generality of certain instructions that require multiple adjustments, making it challenging for agents to achieve full compliance. Inspired by this, we introduce communicative dehallucination, which encourages the assistant to actively seek more detailed suggestions from the instructor before delivering a formal response."
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)
content.alignment = PP_ALIGN.LEFT
content_space = content_frame.add_paragraph()
content_space.text = ""  # Add space between paragraphs

# Ensure text does not exceed the textbox
for paragraph in content_frame.paragraphs:
    paragraph.space_after = Pt(14)

# Add image if needed (example path, replace with actual if provided)
# img_path = "path/to/image.png"
# if os.path.exists(img_path):
#     slide.shapes.add_picture(img_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))