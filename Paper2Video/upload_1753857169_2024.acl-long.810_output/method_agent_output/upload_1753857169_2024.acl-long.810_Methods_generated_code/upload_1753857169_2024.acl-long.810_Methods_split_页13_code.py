import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Minimizing Information Overload"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = ("By sharing only the solutions of each subtask rather than the entire communication history, "
                 "ChatDev minimizes the risk of being overwhelmed by too much information, enhancing concentration "
                 "on each task and encouraging more targeted cooperation, while simultaneously facilitating "
                 "cross-phase context continuity.")
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)
content.alignment = PP_ALIGN.LEFT

# Add a table if needed (not applicable here, but included for completeness)
# table_data = [['Subtask', 'Solution'], ['1', 'Solution 1'], ['2', 'Solution 2']]
# rows, cols = len(table_data), len(table_data[0])
# left = Inches(1)
# top = Inches(3.5)
# width = Inches(8)
# height = Inches(1.5)
# table = slide.shapes.add_table(rows, cols, left, top, width, height).table
# for r in range(rows):
#     for c in range(cols):
#         table.cell(r, c).text = table_data[r][c]
#         table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(16)