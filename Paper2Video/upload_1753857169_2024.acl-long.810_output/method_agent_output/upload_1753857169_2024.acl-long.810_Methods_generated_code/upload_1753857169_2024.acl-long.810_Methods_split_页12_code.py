import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Dialogue Perception through Previous Phases"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "To perceive dialogues through previous phases, the chat chain only transmits the solutions from previous phases as long-term memories tilde { mathcal { M } }, integrating them at the start of the next phase and enabling the cross-phase transmission of long dialogues:"
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Math formula
math_formula = content_frame.add_paragraph()
math_formula.text = "mathcal { T } _ { 1 } ^ { i + 1 } = tilde { mathcal { M } } ^ { i } \\cup mathsf { P } _ { mathcal { T } } ^ { i + 1 }, ~ tilde { mathcal { M } } ^ { i } = bigcup _ { j = 1 } ^ { i } tau ( mathcal { M } _ { | mathcal { M } ^ { j } | } ^ { j } )"
math_formula.font.size = Pt(16)
math_formula.font.color.rgb = RGBColor(0, 0, 0)
math_formula.space_after = Pt(14)

# Additional content
additional_content = content_frame.add_paragraph()
additional_content.text = "where mathsf { P } symbolizes a predetermined prompt that appears exclusively at the start of each phase."
additional_content.font.size = Pt(18)
additional_content.font.color.rgb = RGBColor(0, 0, 0)

# Image (if applicable)
# image_path = "path/to/image.png"  # Replace with actual image path if provided
# if os.path.exists(image_path):
#     slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))