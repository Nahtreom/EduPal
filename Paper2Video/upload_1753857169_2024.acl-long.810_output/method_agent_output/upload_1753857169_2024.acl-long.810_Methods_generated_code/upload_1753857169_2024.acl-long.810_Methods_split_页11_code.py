import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Instruction Generation in Time Step"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "In the next time step t + 1, the instructor utilizes the current memory to generate a new instruction T_{t + 1}^{i}, which is then used to produce a new response A_{t + 1}^{i}. The short-term memory iteratively updates until the number of communications reaches the upper limit |M^{i}|."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Math Formula
formula_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
formula_frame = formula_box.text_frame
formula_frame.word_wrap = True

formula = formula_frame.add_paragraph()
formula.text = "T_{t + 1}^{i} = T(M_{t}^{i}), A_{t + 1}^{i} = A(M_{t}^{i}, T_{t + 1}^{i})"
formula.font.size = Pt(18)
formula.font.color.rgb = RGBColor(0, 0, 0)

formula2 = formula_frame.add_paragraph()
formula2.text = "M_{t + 1}^{i} = M_{t}^{i} ∪ (T_{t + 1}^{i}, A_{t + 1}^{i})"
formula2.font.size = Pt(18)
formula2.font.color.rgb = RGBColor(0, 0, 0)

# Add image if path is provided
image_path = "path/to/image.png"  # Replace with actual image path if available
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5), Inches(4), width=Inches(2), height=Inches(2))