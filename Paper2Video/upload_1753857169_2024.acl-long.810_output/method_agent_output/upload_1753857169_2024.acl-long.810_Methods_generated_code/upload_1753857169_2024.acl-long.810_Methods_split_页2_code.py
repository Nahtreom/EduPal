import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a new slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Challenges in Software Development with LLMs"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = ("Although LLMs show a good understanding of natural and programming languages, "
                 "efficiently transforming textual requirements into functional software in a single step "
                 "remains a significant challenge. ChatDev thus adopts the core principles of the waterfall model, "
                 "using a chat chain (C) with sequential phases (P), each comprising sequential subtasks (T). "
                 "Specifically, ChatDev segments the software development process into three sequential phases: "
                 "design, coding, and testing. The coding phase is further subdivided into subtasks of code writing "
                 "and completion, and the testing phase is segmented into code review (static testing) and system "
                 "testing (dynamic testing), as illustrated in Figure 2.")
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)
content.alignment = PP_ALIGN.LEFT

# Add image if path is provided
image_path = "path/to/your/image.png"  # Replace with actual image path
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(3), height=Inches(2))