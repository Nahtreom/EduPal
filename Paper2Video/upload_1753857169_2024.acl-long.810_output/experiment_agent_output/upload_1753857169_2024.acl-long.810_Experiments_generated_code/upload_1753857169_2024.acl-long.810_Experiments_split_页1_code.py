import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Baselines"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "We chose some representative LLM-based software development methods as our baselines. " \
               "GPT-Engineer (Osika, 2023) is a fundamental single-agent approach in LLM-driven software agents " \
               "with a precise understanding of task requirements and the application of one-step reasoning, " \
               "which highlights its efficiency in generating detailed software solutions at the repository level."
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)
content.space_after = Pt(14)

# Add image if path is provided (example path used here)
image_path = "path/to/image.png"  # Replace with actual image path if available
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5.5), Inches(3), width=Inches(2), height=Inches(2))