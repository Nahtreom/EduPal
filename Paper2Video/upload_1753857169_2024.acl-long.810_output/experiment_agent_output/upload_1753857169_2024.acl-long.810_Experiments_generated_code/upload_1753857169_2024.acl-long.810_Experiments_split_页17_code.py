import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a new slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Task Decomposition and Quality Improvement"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Add first paragraph
p1 = content_frame.add_paragraph()
p1.text = "Therefore, explicitly decomposing the difficult problem into several smaller, more manageable subtasks enhances the effectiveness of task completion."
p1.font.size = Pt(18)
p1.font.color.rgb = RGBColor(0, 0, 0)

# Add second paragraph
p2 = content_frame.add_paragraph()
p2.text = "Additionally, in comparison to MetaGPT, ChatDev significantly raises the Quality from O.1523 to 0.3953."
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0, 0, 0)

# Add a table for comparison (if needed)
# Example: Creating a table for visual comparison
table_data = [
    ['Model', 'Quality'],
    ['MetaGPT', 'O.1523'],
    ['ChatDev', '0.3953']
]

rows, cols = len(table_data), len(table_data[0])
table = slide.shapes.add_table(rows, cols, Inches(5), Inches(2.5), Inches(3), Inches(1.5)).table

# Set column widths
for i in range(cols):
    table.columns[i].width = Inches(1.5)

# Fill the table
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER