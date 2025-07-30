import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a new slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Software Statistics"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "Table 3: Software statistics include Duration (time consumed), #Tokens (number of tokens used), #Files (number of code files generated), and #Lines (total lines of code across all files) in the software generation process."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)
content.space_after = Pt(14)

# Add table
rows, cols = 4, 5
table = slide.shapes.add_table(rows, cols, Inches(1), Inches(4), Inches(8), Inches(2)).table

# Set column names
table.cell(0, 0).text = "Method"
table.cell(0, 1).text = "Duration (s)"
table.cell(0, 2).text = "#Tokens"
table.cell(0, 3).text = "#Files"
table.cell(0, 4).text = "#Lines"

# Fill in the data
data = [
    ["GPT-Engineer", "15.6000", "7,182.5333", "3.9475", "70.2041"],
    ["MetaGPT", "154.0000", "29,278.6510", "4.4233", "153.3000"],
    ["ChatDev", "148.2148", "22,949.4450", "4.3900", "144.3450"]
]

for i, row in enumerate(data, start=1):
    for j, value in enumerate(row):
        table.cell(i, j).text = value

# Set table style
for i in range(rows):
    for j in range(cols):
        cell = table.cell(i, j)
        cell.text_frame.text = cell.text_frame.text.strip()
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER