import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a new slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Understanding User Preferences"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "To further understand user preferences in practical settings, we use the setting adopted by Li et al. (2023a), where agent-generated solutions are compared in pairs by both human participants and the prevalent GPT-4 model to identify the preferred one."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "Table 2 shows ChatDev consistently outperforming other baselines, with higher average win rates in both GPT-4 and human evaluations."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Add a table for comparison (example structure)
table_data = [
    ['Model', 'Average Win Rate'],
    ['ChatDev', '75%'],
    ['Baseline 1', '60%'],
    ['Baseline 2', '55%']
]

rows, cols = len(table_data), len(table_data[0])
left = Inches(5)
top = Inches(3)
width = Inches(3)
height = Inches(1 + 0.5 * (rows - 1))

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
for i in range(cols):
    table.columns[i].width = Inches(1.5)

# Fill the table
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER