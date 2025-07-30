import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Overall Performance"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "As illustrated in Table 1, ChatDev outperforms all baseline methods across all metrics, showing a considerable margin of improvement."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "Firstly, the improvement of ChatDev and MetaGPT over GPT-Engineer demonstrates that complex tasks are difficult to solve in a single-step solution."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Add a table if needed (not included in the original text, but placeholder for future use)
# table_data = [['Metric', 'ChatDev', 'Baseline']]
# table = slide.shapes.add_table(len(table_data), len(table_data[0]), Inches(1), Inches(5), Inches(8), Inches(1)).table
# for i, row in enumerate(table_data):
#     for j, cell in enumerate(row):
#         table.cell(i, j).text = cell
#         table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(16)