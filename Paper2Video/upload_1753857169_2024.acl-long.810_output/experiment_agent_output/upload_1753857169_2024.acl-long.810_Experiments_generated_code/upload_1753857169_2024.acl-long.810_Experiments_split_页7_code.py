import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Overall Performance of LLM-Powered Software Development Methods"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "Table 1: Overall performance of the LLM-powered software development methods, encompassing both single-agent (trianglerighteq { overbrace { 2 / 3 } }) and multi-agent (hat{AA} / z, hat{AA}) paradigms. Performance metrics are averaged for all tasks. The top scores are in bold, with second-highest underlined. † indicates significant statistical differences (p ≤ 0.05) between a baseline and ours."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)
content.alignment = PP_ALIGN.LEFT

# Add a table (if needed, based on the context)
# Example: Creating a table for performance metrics
table_data = [
    ['Method', 'Performance'],
    ['Single-Agent', 'Value1'],
    ['Multi-Agent', 'Value2']
]

rows, cols = len(table_data), len(table_data[0])
table = slide.shapes.add_table(rows, cols, Inches(1), Inches(5.5), Inches(8), Inches(1)).table

# Set column widths
for i in range(cols):
    table.columns[i].width = Inches(4)

# Fill the table
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER