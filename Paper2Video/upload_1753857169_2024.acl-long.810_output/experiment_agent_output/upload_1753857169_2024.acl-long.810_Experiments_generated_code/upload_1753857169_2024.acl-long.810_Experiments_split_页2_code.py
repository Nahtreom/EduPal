import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "MetaGPT Framework"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "MetaGPT (Hong et al., 2023) is an advanced framework that allocates specific roles to various LLM-driven software agents and incorporates standardized operating procedures to enable multi-agent participation. In each step, agents with specific roles generate solutions by adhering to static instructions predefined by human experts."
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)
content.alignment = PP_ALIGN.LEFT

# Add image if path is provided (example path)
image_path = "path/to/image.png"  # Replace with actual image path if available
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))