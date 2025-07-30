import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a new slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Dataset Overview"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "This dataset includes important software categories from popular platforms such as Ubuntu, Google Play, Microsoft Store, and Apple Store."
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "It comprises 1,200 software task prompts that have been carefully categorized into 5 main areas: Education, Work, Life, Game, and Creation."
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "All these areas are further divided into 40 subcategories, and each subcategory contains 30 unique task prompts."
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)

# Add a table if needed (not specified in the text, so this is optional)
# Example of adding a table can be included here if data comparison is provided