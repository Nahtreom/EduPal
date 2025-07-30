import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a new slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Comparison of Methods"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "Method: GPT-Engineer",
    "Completeness: 0.5022t",
    "Executability: 0.3583t",
    "Consistency: 0.7887+",
    "Quality: 0.1419t",
    "",
    "Method: MetaGPT",
    "Completeness: 0.4834†",
    "Executability: 0.4145†",
    "Consistency: 0.7601t",
    "Quality: 0.1523t",
    "",
    "Method: ChatDev",
    "Completeness: 0",
    "Executability: 0.5600",
    "Consistency: 0.8800",
    "Quality: 0.3953"
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(10)

# Add table
table_data = [
    ["Method", "Paradigm", "Completeness", "Executability", "Consistency", "Quality"],
    ["GPT-Engineer", "", "0.5022t", "0.3583t", "0.7887+", "0.1419t"],
    ["MetaGPT", "恩", "0.4834†", "0.4145†", "0.7601t", "0.1523t"],
    ["ChatDev", "0", "0.5600", "0.8800", "0.8021", "0.3953"]
]

rows, cols = len(table_data), len(table_data[0])
left = Inches(5)
top = Inches(2)
width = Inches(3)
height = Inches(2)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
for i in range(cols):
    table.columns[i].width = Inches(0.5)

# Fill the table
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER