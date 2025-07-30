import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Metrics"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "Evaluating software is also a challenging task, especially when trying to assess it on a holistic level. Under the current limitation of scarce benchmark resources, traditional function-oriented code generation metrics (e.g., pass@k) cannot seamlessly transfer to a comprehensive evaluation of entire software systems."
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "The main reason for this is that it is often impractical to develop manual or automated test cases for various types of software, especially those involving complex interfaces, frequent user interactions, or non-deterministic feedback."
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)

# Add image if needed (example path, replace with actual path if available)
# img_path = "path/to/image.png"
# if os.path.exists(img_path):
#     slide.shapes.add_picture(img_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))