import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Implementation Details"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "We divided software development into 5 subtasks within 3 phases, assigning specific roles like CEO, CTO, programmer, reviewer, and tester."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "A subtask would terminate and get a conclusion either after two unchanged code modifications or after 10 rounds of communication."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "During the code completion, review, and testing, a communicative dehallucination is activated."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Add a placeholder for an image (if needed)
# image_path = "path_to_image"  # Replace with actual image path if available
# if os.path.exists(image_path):
#     slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))