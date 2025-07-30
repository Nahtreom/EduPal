import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Executability Assessment"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = "Executability assesses the software's ability to run correctly within a compilation environment, quantified as the percentage of software that compiles successfully and can run directly. A higher score indicates a higher probability of successful execution."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Adjust alignment
for paragraph in content_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT

# Add a table if needed (not applicable here, but structure provided)
# table_data = [['Metric', 'Value'], ['Executability Score', '85%']]
# rows, cols = len(table_data), len(table_data[0])
# table = slide.shapes.add_table(rows, cols, Inches(1), Inches(5.5), Inches(8), Inches(1)).table
# for r in range(rows):
#     for c in range(cols):
#         table.cell(r, c).text = table_data[r][c]
#         table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(16)