import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a new slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Solution Identification"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content = content_frame.add_paragraph()
content.text = "For ease of identifying solutions, the assistant begins responses with \"<SOLUTION>\" when a consensus is reached."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "We used ChatGPT-3.5 with a temperature of 0.2 and integrated Python-3.11.4 for feedback."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

content = content_frame.add_paragraph()
content.text = "All baselines in the evaluation share the same hyperparameters and settings for fairness."
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)

# Adjust text box properties
for paragraph in content_frame.paragraphs:
    paragraph.space_after = Pt(14)
    paragraph.alignment = PP_ALIGN.LEFT

# Add a table if needed (not applicable in this case)