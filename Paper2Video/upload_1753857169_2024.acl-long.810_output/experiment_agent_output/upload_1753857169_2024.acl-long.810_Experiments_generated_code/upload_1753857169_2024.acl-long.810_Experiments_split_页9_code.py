import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a new slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Table 2: Pairwise evaluation results."
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Add table
table_data = [
    ["Method", "Evaluator", "Baseline Wins", "ChatDev Wins", "Draw"],
    ["GPT-Engineer", "", "", "", ""],
    ["", "GPT-4", "22.50%", "77.08%", "00.42%"],
    ["MetaGPT", "GPT-4", "37.50%", "57.08%", "05.42%"],
    ["", "Human", "07.92%", "88.00%", "04.08%"]
]

rows, cols = len(table_data), len(table_data[0])
left = Inches(1)
top = Inches(4.5)
width = Inches(8)
height = Inches(1 + 0.5 * (rows - 1))

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
for i in range(cols):
    table.columns[i].width = Inches(1.5)

# Fill the table
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(220, 220, 220) if r == 0 else RGBColor(255, 255, 255)
        cell.border_left.fill.solid()
        cell.border_left.fill.fore_color.rgb = RGBColor(0, 0, 0)
        cell.border_top.fill.solid()
        cell.border_top.fill.fore_color.rgb = RGBColor(0, 0, 0)
        cell.border_right.fill.solid()
        cell.border_right.fill.fore_color.rgb = RGBColor(0, 0, 0)
        cell.border_bottom.fill.solid()
        cell.border_bottom.fill.fore_color.rgb = RGBColor(0, 0, 0)