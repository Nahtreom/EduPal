import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Technical Inconsistencies in Development Phases"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "Due to these technical inconsistencies, methods employed in different phases remain isolated until now.",
    "Every phase, from data collection and labeling to model training and inference, requires its unique designs,",
    "leading to a fragmented and less efficient development process in the field (Freeman et al., 2001; Ernst, 2017; Winkler et al., 2020).",
    "Motivated by the expert-like potential of autonomous agents, we aim to establish language as a unifying bridgeâ€”",
    "utilizing multiple LLM-powered agents with specialized roles for cooperative software development through language-based communication across different phases;",
    "solutions in different phases are derived from their multi-turn dialogues, whether dealing with text or code."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)

# Add image if path is provided
image_path = "path/to/image.png"  # Replace with actual image path if available
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(5.5), Inches(3), width=Inches(2), height=Inches(2))