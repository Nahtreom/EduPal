import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "ChatDev's Dehallucination Mechanism"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = ("Additionally, to minimize coding hallucinations, ChatDev includes a communicative dehallucination mechanism, "
                 "enabling agents to actively request more specific details before giving direct responses. The communication pattern "
                 "instructs agents on how to communicate, enabling precise information exchange for effective solution optimization "
                 "while reducing coding hallucinations. We built a comprehensive dataset containing software requirement descriptions "
                 "and conducted comprehensive analyses. The results indicate that ChatDev notably improves the quality of software, "
                 "leading to improved completeness, executability, and better consistency with requirements.")
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)
content.alignment = PP_ALIGN.LEFT

# Add a table for data comparison (if applicable)
# Example: Creating a table for hypothetical data comparison
table_data = [
    ['Metric', 'Before ChatDev', 'After ChatDev'],
    ['Completeness', '75%', '90%'],
    ['Executability', '70%', '85%'],
    ['Consistency', '65%', '88%']
]

rows, cols = len(table_data), len(table_data[0])
left = Inches(5)
top = Inches(2)
width = Inches(3)
height = Inches(2)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
for i in range(cols):
    table.columns[i].width = Inches(1)

# Fill the table with data
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER