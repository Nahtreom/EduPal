import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Challenges of Coding Hallucinations"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = ("Nevertheless, due to the tendency of LLM hallucinations (Dhuliawala et al., 2023; "
                 "Zhang et al., 2023b), the strategy of generating software through communicative agents "
                 "could lead to the non-trivial challenge of coding hallucinations, which involves the "
                 "generation of source code that is incomplete, unexecutable, or inaccurate, ultimately "
                 "failing to fulfill the intended requirements (Agnihotri and Chug, 2020). The frequent "
                 "occurrence of coding hallucination in turn reflects the constrained autonomy of agents "
                 "in task completion, inevitably demanding additional manual intervention and thereby "
                 "hindering the immediate usability and reliability of the generated software (Ji et al., 2023).")
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)
content.alignment = PP_ALIGN.LEFT

# Optional: Add an image if a path is provided
# image_path = "path_to_image"  # Replace with actual image path if available
# if os.path.exists(image_path):
#     slide.shapes.add_picture(image_path, Inches(6), Inches(3), width=Inches(2), height=Inches(2))