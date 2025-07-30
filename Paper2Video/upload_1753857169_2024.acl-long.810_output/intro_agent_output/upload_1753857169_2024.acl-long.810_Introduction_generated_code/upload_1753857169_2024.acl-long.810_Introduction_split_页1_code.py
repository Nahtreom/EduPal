import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Create a slide with a blank layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Transformations in Large Language Models"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "Large language models (LLMs) have led to substantial transformations due to their ability to effortlessly integrate extensive knowledge expressed in language (Brown et al., 2020; Bubeck et al., 2023), combined with their strong capacity for roleplaying within designated roles (Park et al., 2023; Hua et al., 2023; Chen et al., 2023b).",
    "This advancement eliminates the need for model-specific designs and delivers impressive performance in diverse downstream applications.",
    "Furthermore, autonomous agents (Richards, 2023; Zhou et al., 2023a) have gained attention for enhancing the capabilities of LLMs with advanced features such as context-aware memory (Sumers et al., 2023), multistep planning (Liu et al., 2023), and strategic tool using (Schick et al., 2023)."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)

# Add image if path is provided
# image_path = "path/to/image.png"  # Replace with actual image path if available
# if os.path.exists(image_path):
#     slide.shapes.add_picture(image_path, Inches(5), Inches(3), width=Inches(2), height=Inches(2))