import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "ChatDev: A Chat-Powered Software-Development Framework"
title.font.size = Pt(24)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "In this paper, we propose ChatDev, a chat-powered software-development framework integrating multiple 'software agents' for active involvement in three core phases of the software lifecycle: design, coding, and testing.",
    "Technically, ChatDev uses a chat chain to divide each phase into smaller subtasks further, enabling agents' multi-turn communications to cooperatively propose and develop solutions (e.g., creative ideas or source code).",
    "The chain-structured workflow guides agents on what to communicate, fostering cooperation and smoothly linking natural- and programming-language subtasks to propel problem solving."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)

# Image
img_path = "path/to/your/image.png"  # Replace with actual image path
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(5), Inches(3), width=Inches(3), height=Inches(2))