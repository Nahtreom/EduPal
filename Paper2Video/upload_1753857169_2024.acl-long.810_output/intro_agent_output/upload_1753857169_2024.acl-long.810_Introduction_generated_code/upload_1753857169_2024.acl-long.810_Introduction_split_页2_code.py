import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Software Development Complexity"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = [
    "Software development is a complex task that necessitates cooperation among multiple members with diverse skills (e.g., architects, programmers, and testers) (Basili, 1989; Sawyer and Guinan, 1998).",
    "This entails extensive communication among different roles to understand and analyze requirements through natural language, while also encompassing development and debugging using programming languages (Ernst, 2017; Banker et al., 1998).",
    "Numerous studies use deep learning to improve specific phases of the waterfall model in software development, such as design, coding, and testing (Pudlitz et al., 2019; Martin and Abran, 2015; Gao et al., 2019; Wang et al., 2016)."
]

for line in content:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)

# Optional: Add a table for data comparison if needed
# table_data = [['Phase', 'Deep Learning Impact'], ['Design', 'Improved efficiency'], ['Coding', 'Error reduction'], ['Testing', 'Enhanced accuracy']]
# rows, cols = len(table_data), len(table_data[0])
# left = Inches(1)
# top = Inches(5.5)
# width = Inches(8)
# height = Inches(1.5)
# table = slide.shapes.add_table(rows, cols, left, top, width, height).table
# for r in range(rows):
#     for c in range(cols):
#         table.cell(r, c).text = table_data[r][c]
#         table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(16)