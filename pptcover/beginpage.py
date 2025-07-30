from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 替换变量
title_text = "(可替换)标题信息"
author_text = "(可替换)作者信息"
affiliation_text = "(可替换)单位信息, 例如：大学A, 研究所B"

# 创建演示文稿
#prs = Presentation()

# 设置幻灯片尺寸（可选）
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 添加一张空白幻灯片
slide_layout = prs.slide_layouts[6]  # 空白布局
slide = prs.slides.add_slide(slide_layout)

# 添加背景图，需要替换为实际图片路径
background_image_path = "/home/EduAgent/pptcover/cover.png"
pic = slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
# 调整图片层级到最底层
pic.z_order = 0

# 添加文本框，设置位置和大小（让它尽量居中展示）
left = Inches(1)
top = Inches(2)
width = prs.slide_width - Inches(2)
height = Inches(3.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True  # ✅ 启用自动换行
text_frame.clear()  # 清除默认段落

# 添加标题段落
p1 = text_frame.add_paragraph()
p1.text = title_text
p1.font.size = Pt(40)
p1.font.color.rgb = RGBColor(0, 0, 0)
p1.alignment = PP_ALIGN.CENTER

# 添加作者段落
p2 = text_frame.add_paragraph()
p2.text = author_text
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(0, 0, 0)
p2.alignment = PP_ALIGN.CENTER

# 添加单位段落
p3 = text_frame.add_paragraph()
p3.text = affiliation_text
p3.font.size = Pt(24)
p3.font.color.rgb = RGBColor(0, 0, 0)
p3.alignment = PP_ALIGN.CENTER

# 保存PPT
#prs.save("cover_slide.pptx")
