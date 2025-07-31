import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# 创建Presentation对象
#prs = Presentation()

# 使用一个空白布局（白底）
slide_layout = prs.slide_layouts[6]  # 空白布局
slide = prs.slides.add_slide(slide_layout)

# 插入照片的路径，需要替换为实际的图片路径
#image_path = '/home/EduAgent/pptcover/conclusion.jpg'

# 获取幻灯片的宽度和高度
slide_width = prs.slide_width
slide_height = prs.slide_height

# 设置图片位置为幻灯片左上角，大小为幻灯片的宽高
left = top = 0
width = slide_width
height = slide_height


title_text = "Experiments"

left = Inches(1)
top = Inches(2)
width = prs.slide_width - Inches(2)
height = Inches(3.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True  # ✅ 启用自动换行
text_frame.clear()  # 清除默认段落

# 添加标题段落
p1 = text_frame.add_paragraph()
p1.text = title_text
p1.font.size = Pt(40)
p1.font.color.rgb = RGBColor(0, 0, 0)
p1.alignment = PP_ALIGN.CENTER

# 在幻灯片中插入图片
#slide.shapes.add_picture(image_path, left, top, width=width, height=height)
