#!/home/EduAgent/miniconda3/envs/edu_env/bin/python3
# -*- coding: utf-8 -*-

import re
import os
import glob
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData


# 从文件夹中找到所有 .py 文件
def find_ppt_code_files(code_dir):
    """查找指定目录下所有的 .py 文件"""
    return glob.glob(os.path.join(code_dir, "*.py"))

def create_slide_from_code(prs, slide_code):
    """从给定的PPT代码字符串创建一页PPT并添加到PPT中"""
    try:
        # 使用exec来执行PPT代码，并确保prs变量被正确传递和修改
        exec(slide_code, {"prs": prs, "__name__": "__main__"})
    except Exception as e:
        print(f"⚠️ 执行PPT代码时出错: {e}")

# def generate_complete_ppt_from_folder(folder_path, output_ppt_path):
#     """从文件夹中的所有 .py 文件生成一个完整的PPT"""
#     prs = Presentation()  # 创建一个新的PPT对象

#     # 获取文件夹中的所有 .py 文件
#     ppt_files = find_ppt_code_files(folder_path)
    
#     if not ppt_files:
#         print("❌ 未找到任何PPT代码文件。")
#         return
    
#     print(f"📊 找到 {len(ppt_files)} 个PPT代码文件，正在生成PPT...")

#     for ppt_file in ppt_files:
#         print(f"📄 正在处理: {os.path.basename(ppt_file)}")
        
#         # 读取每个PPT代码文件
#         with open(ppt_file, "r", encoding="utf-8") as f:
#             slide_code = f.read()

#         # 使用 exec 执行PPT代码，并将生成的页面添加到PPT中
#         create_slide_from_code(prs, slide_code)

#     # 保存完整PPT
#     prs.save(output_ppt_path)
#     print(f"🎉 完成！完整的PPT已保存到: {output_ppt_path}")
def generate_complete_ppt_from_folder(folder_path, output_ppt_path):
    """从文件夹中的所有 .py 文件生成一个完整的PPT"""
    prs = Presentation()  # 创建一个新的PPT对象
    
    prs.slide_width = Inches(12)
    prs.slide_height = Inches(8)

    # 获取文件夹中的所有 .py 文件
    ppt_files = find_ppt_code_files(folder_path)
    
    if not ppt_files:
        print("❌ 未找到任何PPT代码文件。")
        return
    
    print(f"📊 找到 {len(ppt_files)} 个PPT代码文件，正在生成PPT...")

    # 定义类别顺序
    categories = ['pptcover','Introduction', 'Methods', 'Experiments', 'Conclusion']

    ppt_files = [
    f for f in ppt_files
    if re.search(r'(pptcover|Introduction|Methods|Experiments|Conclusion)', os.path.basename(f))
]

    # 自定义排序规则：首先按类别排序，然后按页号排序
    ppt_files.sort(key=lambda f: (categories.index(re.search(r'(pptcover|Introduction|Methods|Experiments|Conclusion)', os.path.basename(f)).group(0)),
                                  int(re.search(r'页(\d+)', os.path.basename(f)).group(1)) if re.search(r'页(\d+)', os.path.basename(f)) else 0))

    for ppt_file in ppt_files:
        print(f"📄 正在处理: {os.path.basename(ppt_file)}")
        
        # 读取每个PPT代码文件
        with open(ppt_file, "r", encoding="utf-8") as f:
            slide_code = f.read()

        # 使用 exec 执行PPT代码，并将生成的页面添加到PPT中
        create_slide_from_code(prs, slide_code)
    
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 获取幻灯片的宽度和高度
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 设置图片位置为幻灯片左上角，大小为幻灯片的宽高
    left = top = 0
    width = slide_width
    height = slide_height


    title_text = "Thanks"

    left = Inches(1)
    top = Inches(2)
    width = slide_width - Inches(2)
    height = Inches(3.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True  # ✅ 启用自动换行
    text_frame.clear()  # 清除默认段落

    # 添加标题段落
    p1 = text_frame.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(46)
    p1.font.color.rgb = RGBColor(0, 0, 0)
    p1.alignment = PP_ALIGN.CENTER

    # 保存完整PPT
    prs.save(output_ppt_path)
    print(f"🎉 完成！完整的PPT已保存到: {output_ppt_path}")

def main():
    parser = argparse.ArgumentParser(
        description="根据指定的输出目录生成一个完整的PPT，合并所有的Python PPT代码"
    )
    parser.add_argument(
        "output_dir",
        help="指定子目录名称：脚本会去 Paper2Video/<output_dir>/final_results/Code 查找 .py 文件并输出到最终的PPT"
    )
    args = parser.parse_args()

    # 获取脚本所在目录
    script_dir = os.path.dirname(os.path.abspath(__file__))
    base_dir = os.path.join(script_dir, "Paper2Video", args.output_dir, "final_results")
    code_dir = os.path.join(base_dir, "Code")  # 获取 Code 目录路径

    if not os.path.isdir(code_dir):
        print(f"❌ 未找到代码目录: {code_dir}")
        return

    output_ppt_path = os.path.join(base_dir, "final_presentation.pptx")

    # 生成完整PPT
    generate_complete_ppt_from_folder(code_dir, output_ppt_path)

if __name__ == "__main__":
    main()
